import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "SISTEM INFORMASI SUMBER DAYA MANUSIA (HRIS)"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 0, 0)
    
    subtitle.text = "PT DEA GLOBAL NIAGA\nSolusi Otomatisasi & Digitalisasi HRD"

    # Slide 1: Permasalahan (Masalah)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Latar Belakang & Permasalahan"
    body = slide.placeholders[1].text_frame
    body.text = "Proses HR Manual Menimbulkan Tantangan Besar:"
    p = body.add_paragraph()
    p.text = "1. Pencatatan Jam Kerja Tidak Akurat (Risiko Manipulasi/Titip Absen)."
    p.level = 1
    p = body.add_paragraph()
    p.text = "2. Kehilangan Data Cuti (Pengajuan berbasis kertas/pesan singkat)."
    p.level = 1
    p = body.add_paragraph()
    p.text = "3. Efisiensi Waktu Terbuang (Perekapan gaji & absensi bulanan lambat)."
    p.level = 1
    p = body.add_paragraph()
    p.text = "4. Transparansi Kinerja Rendah."
    p.level = 1

    # Slide 2: Solusi Utama (Solusi)
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Solusi HRIS Berbasis Web"
    body = slide.placeholders[1].text_frame
    body.text = "Transformasi Digital melalui HRIS Modern:"
    p = body.add_paragraph()
    p.text = "Kehadiran Berbasis Lokasi (GPS) & Pengenalan Wajah (Face API)."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Pengajuan Cuti & Izin Digital secara Real-time."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Sistem PWA (Progressive Web App): Bisa di-Install di Android & iOS."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Dashboard Analitik Lanjutan bagi Manajemen."
    p.level = 1

    # Slide 3: Fitur 1 - Autentikasi
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Fitur: Keamanan & Autentikasi"
    body = slide.placeholders[1].text_frame
    body.text = "Pemisahan Hak Akses Multi-Level:"
    p = body.add_paragraph()
    p.text = "Admin/HR: Memiliki akses penuh ke pelaporan dan persetujuan."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Pegawai: Hanya dapat melihat data pribadi dan melakukan presensi."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Keamanan Password dengan Bcrypt (Hash 10-rounds)."
    p.level = 1
    
    if os.path.exists('diagrams/auth_flow.png'):
        slide.shapes.add_picture('diagrams/auth_flow.png', Inches(5.5), Inches(1.5), width=Inches(4.0))

    # Slide 4: Fitur 2 - Presensi Pintar
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Fitur: Presensi Pintar (Smart Attendance)"
    body = slide.placeholders[1].text_frame
    body.text = "Meniadakan Kecurangan dengan AI:"
    p = body.add_paragraph()
    p.text = "Deteksi Wajah Karyawan secara langsung dari browser."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Validasi Koordinat GPS saat menekan tombol Check-In."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Watermark otomatis pada foto (Waktu & Lokasi)."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Kompresi Gambar Otomatis untuk menghemat penyimpanan cloud."
    p.level = 1
    
    if os.path.exists('screenshots/3_attendance.png'):
        slide.shapes.add_picture('screenshots/3_attendance.png', Inches(5.0), Inches(1.5), width=Inches(4.5))

    # Slide 5: PWA & Kemudahan Akses
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Kemudahan Instalasi: Progressive Web App (PWA)"
    body = slide.placeholders[1].text_frame
    body.text = "Tidak Perlu App Store atau Play Store:"
    p = body.add_paragraph()
    p.text = "Buka link Vercel di browser (Safari/Chrome)."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Pilih 'Tambahkan ke Layar Utama' (Add to Home Screen)."
    p.level = 1
    p = body.add_paragraph()
    p.text = "Tampil sebagai Aplikasi Native, lebih cepat & hemat kuota."
    p.level = 1

    # Slide 6: Rancangan Sistem (ERD)
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Perancangan Sistem: ERD"
    
    if os.path.exists('diagrams/erd.png'):
        slide.shapes.add_picture('diagrams/erd.png', Inches(1.5), Inches(1.5), width=Inches(7.0))

    # Slide 7: Rancangan Sistem (Use Case)
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Perancangan Sistem: Use Case Diagram"
    
    if os.path.exists('diagrams/use_case.png'):
        slide.shapes.add_picture('diagrams/use_case.png', Inches(2.0), Inches(1.5), width=Inches(6.0))

    # Slide 8: Modul Cuti & Manajemen (Employees, Reports, Profile)
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Modul Cuti, Laporan, & Manajemen Karyawan"
    
    if os.path.exists('screenshots/5_employees.png'):
        slide.shapes.add_picture('screenshots/5_employees.png', Inches(1.0), Inches(1.5), width=Inches(3.5))
    if os.path.exists('screenshots/6_reports.png'):
        slide.shapes.add_picture('screenshots/6_reports.png', Inches(5.0), Inches(1.5), width=Inches(3.5))

    # Slide 9: Penutup & Kesimpulan
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Kesimpulan & ROI"
    
    body = slide.shapes.placeholders[1].text_frame
    p = body.add_paragraph()
    p.text = "Sistem telah memodernisasi cara perusahaan melacak absensi."
    p.level = 0
    p = body.add_paragraph()
    p.text = "Meningkatkan kedisiplinan dan mencegah fraud dengan teknologi cerdas."
    p.level = 0
    p = body.add_paragraph()
    p.text = "Menghemat waktu ratusan jam per bulan untuk perekapan gaji."
    p.level = 0

    prs.save('Presentasi_HRIS.pptx')
    print("Presentasi PPTX berhasil dibuat!")

if __name__ == '__main__':
    create_presentation()
